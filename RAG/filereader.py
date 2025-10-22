import os, fitz, docx, pandas as pd
from pptx import Presentation
from pathlib import Path

# ---------- File-type handlers ----------

def read_txt(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()

def read_pdf(path):
    doc = fitz.open(path)
    pages = [page.get_text("text") for page in doc]
    doc.close()
    return "\n".join(pages)

def read_docx(path):
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs if p.text.strip())

def read_pptx(path):
    prs = Presentation(path)
    slides=[]
    for i,slide in enumerate(prs.slides,1):
        texts=[s.text for s in slide.shapes if hasattr(s,"text") and s.text.strip()]
        slides.append(f"--- Slide {i} ---\n"+"\n".join(texts))
    return "\n".join(slides)

def read_csv_xlsx(path):
    try:
        df = pd.read_csv(path) if path.endswith(".csv") else pd.read_excel(path)
        return df.to_string(index=False)
    except Exception as e:
        return f"⚠️ Could not read table: {e}"

# ---------- Main unified reader ----------

def extract_text(file_path:str)->str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":   return read_pdf(file_path)
    if ext == ".docx":  return read_docx(file_path)
    if ext == ".pptx":  return read_pptx(file_path)
    if ext in (".csv",".xlsx"): return read_csv_xlsx(file_path)
    if ext == ".txt":   return read_txt(file_path)
    return "⚠️ Unsupported file type."

if __name__ == "__main__":
    path = input("Enter file path: ").strip().strip('"')
    if not os.path.exists(path):
        print("❌ File not found.")
    else:
        print("\n🧠 Extracting text...\n")
        text = extract_text(path)
        print(text[:1500])  # print first 1500 chars
        print("\n✅ Extraction complete.")
